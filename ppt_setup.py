#Import Presentation from pptx module (pip install module)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt




#Create PPT Presentation 
prs = Presentation()

#Assigning layout to slide 1 - Title and Empty space
slide_layout = prs.slide_layouts[6]

#Adding slide 1 to presentation
slide = prs.slides.add_slide(slide_layout)

background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(180, 206, 228)

#Adding text box to appear as title
left = Inches(2.5)
top = Inches(2)
width = height = Inches(1)
txBox = slide.shapes.add_textbox(left, top, width, height)
t_in = txBox.text_frame

p = t_in.add_paragraph()
p.text = 'print("Python")'
p.font.bold = True
p.font.size = Pt(50)
p.font.color.rgb = RGBColor(0, 0, 0)


#Adding text box to appear in the middle of the slide, below title
left1 = Inches(1.5)
top1 = Inches(4)
width1 = height1 = Inches(1)
txBox1 = slide.shapes.add_textbox(left1, top1, width1, height1)
t_in1 = txBox1.text_frame

p1 = t_in1.add_paragraph()
p1.text = "For those seeking a simpler way of life..."
p1.font.bold = True
p1.font.size = Pt(30)
p1.font.color.rgb = RGBColor(0, 0, 0)


###############################################################################

#Assigning layout to slide 2 - Title and Content
slide2_layout = prs.slide_layouts[1]

#Adding slide 2 to presentation
slide2 = prs.slides.add_slide(slide2_layout)
shapes = slide2.shapes

background = slide2.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(253, 239, 180)

#Slide 2 title placeholder
title2 = slide2.shapes.title
b_shape = shapes.placeholders[1]

#Adding text to slide 2 title placeholder
title2.text = "What is Python?"

tf = b_shape.text_frame

#Adding bullet point 1
tf1 = tf.add_paragraph()
tf1.text = "Released in 1991"
tf1.font.size = Pt(25)
tf1.level = 0

#Adding bullet point 2
tf2 = tf.add_paragraph()
tf2.text = "Python is an object-oriented programming language (OOP)"
tf2.font.size = Pt(25)
tf2.level = 0

#Adding bullet point 3
tf3 = tf.add_paragraph()
tf3.text = "Beginner-friendly due to its simpler syntax - write programs with fewer lines of code"
tf3.font.size = Pt(25)
tf3.level = 0

#Adding bullet point 4
tf4 = tf.add_paragraph()
tf4.text = "Not only used by software engineers - also by mathmatecians, data analysts, scientists, accountants and network engineers"
tf4.font.size = Pt(25)
tf4.level = 0

#Adding bullet point 5
tf5 = tf.add_paragraph()
tf5.text = "One of the most popular programming languages in the world!"
tf5.font.size = Pt(25)
tf5.level = 0

###############################################################################

#Assigning layout to slide 3 - Title and Empty space
slide3_layout = prs.slide_layouts[5]

#Adding slide 3 to presentation
slide3 = prs.slides.add_slide(slide3_layout)

background = slide3.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(180, 206, 228)

#Slide 3 title placeholder
title3 = slide3.shapes.title

#Adding text to slide 3 title
title3.text = "Are you ready to RUUUMBLE?! - Python vs JavaScript"

#Adding text box to appear in the middle of the slide, below title
left3 = Inches(3)
top3 = Inches(2)
width3 = height3 = Inches(1)
txBox3 = slide3.shapes.add_textbox(left3, top3, width3, height3)
t_in3 = txBox3.text_frame

p3 = t_in3.add_paragraph()
p3.text = "Naming Conventions"
p3.font.bold = True
p3.font.size = Pt(30)

#Adding images
img_naming_js = "img/naming_js.PNG"
img_naming_py = "img/naming_py.PNG"

from_left_naming_js = Inches(1.5)
from_top_naming_js = Inches(4)
from_left_naming_py = Inches(5.5)
from_top_naming_py = Inches(4)
add_naming_js = slide3.shapes.add_picture(img_naming_js, from_left_naming_js, from_top_naming_js)
add_naming_py = slide3.shapes.add_picture(img_naming_py, from_left_naming_py, from_top_naming_py)

###############################################################################

#Assigning layout to slide 4 - Title and Empty space
slide4_layout = prs.slide_layouts[5]

#Adding slide 4 to presentation
slide4 = prs.slides.add_slide(slide4_layout)

background = slide4.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(253, 239, 180)

#Slide 4 title placeholder
title4 = slide4.shapes.title

#Adding text to slide 4 title
title4.text = "Are you ready to RUUUMBLE?! - Python vs JavaScript"

#Adding text box to appear in the middle of the slide, below title
left4 = Inches(3.7)
top4 = Inches(2)
width4 = height4 = Inches(1)
txBox4 = slide4.shapes.add_textbox(left4, top4, width4, height4)
t_in4 = txBox4.text_frame

p4 = t_in4.add_paragraph()
p4.text = "Code Blocks"
p4.font.bold = True
p4.font.size = Pt(30)

#Adding images
img_code_js = "img/code_block_js.PNG"
img_code_py = "img/code_block_py.PNG"

from_left_code_js = Inches(2)
from_top_code_js = Inches(3)
from_left_code_py = Inches(2.3)
from_top_code_py = Inches(5.5)
add_code_js = slide4.shapes.add_picture(img_code_js, from_left_code_js, from_top_code_js)
add_code_py = slide4.shapes.add_picture(img_code_py, from_left_code_py, from_top_code_py)


###############################################################################

#Assigning layout to slide 5 - Title and Empty space
slide5_layout = prs.slide_layouts[5]

#Adding slide 5 to presentation
slide5 = prs.slides.add_slide(slide5_layout)

background = slide5.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(180, 206, 228)

#Slide 5 title placeholder
title5 = slide5.shapes.title

#Adding text to slide 5 title
title5.text = "Are you ready to RUUUMBLE?! - Python vs JavaScript"

#Adding text box to appear in the middle of the slide, below title
left5 = Inches(3.3)
top5 = Inches(2)
width5 = height5 = Inches(1)
txBox5 = slide5.shapes.add_textbox(left5, top5, width5, height5)
t_in5 = txBox5.text_frame

p5 = t_in5.add_paragraph()
p5.text = "Defining Variables"
p5.font.bold = True
p5.font.size = Pt(30)

#Adding images
img_var_js = "img/variables_js.PNG"
img_var_py = "img/variables_py.PNG"

from_left_var_js = Inches(2.7)
from_top_var_js = Inches(3.5)
from_left_var_py = Inches(3.2)
from_top_var_py = Inches(5)
add_var_js = slide5.shapes.add_picture(img_var_js, from_left_var_js, from_top_var_js)
add_var_py = slide5.shapes.add_picture(img_var_py, from_left_var_py, from_top_var_py)

###############################################################################

#Assigning layout to slide 6 - Title and Empty space
slide6_layout = prs.slide_layouts[5]

#Adding slide 6 to presentation
slide6 = prs.slides.add_slide(slide6_layout)

background = slide6.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(253, 239, 180)

#Slide 6 title placeholder
title6 = slide6.shapes.title

#Adding text to slide 6 title
title6.text = "Are you ready to RUUUMBLE?! - Python vs JavaScript"

#Adding text box to appear in the middle of the slide, below title
left6 = Inches(3.7)
top6 = Inches(2)
width6 = height6 = Inches(1)
txBox6 = slide6.shapes.add_textbox(left6, top6, width6, height6)
t_in6 = txBox6.text_frame

p6 = t_in6.add_paragraph()
p6.text = "FOR Loops"
p6.font.bold = True
p6.font.size = Pt(30)

#Adding images
img_loop_js = "img/for_loop_js.PNG"
img_loop_py = "img/for_loop_py.PNG"

from_left_loop_js = Inches(1.2)
from_top_loop_js = Inches(3.3)
from_left_loop_py = Inches(5.5)
from_top_loop_py = Inches(3.3)
add_loop_js = slide6.shapes.add_picture(img_loop_js, from_left_loop_js, from_top_loop_js)
add_loop_py = slide6.shapes.add_picture(img_loop_py, from_left_loop_py, from_top_loop_py)


###############################################################################

#Assigning layout to slide 7 - Title and Empty space
slide7_layout = prs.slide_layouts[5]

#Adding slide 7 to presentation
slide7 = prs.slides.add_slide(slide7_layout)

background = slide7.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(180, 206, 228)

#Slide 7 title placeholder
title7 = slide7.shapes.title

#Adding text to slide 7 title
title7.text = "Are you ready to RUUUMBLE?! - Python vs JavaScript"

#Adding text box to appear in the middle of the slide, below title
left7 = Inches(3.3)
top7 = Inches(2)
width7 = height7 = Inches(1)
txBox7 = slide7.shapes.add_textbox(left7, top7, width7, height7)
t_in7 = txBox7.text_frame

p7 = t_in7.add_paragraph()
p7.text = "if/else Statements"
p7.font.bold = True
p7.font.size = Pt(30)

#Adding images
img_if_js = "img/if_else_js.PNG"
img_if_py = "img/if_else_py.PNG"

from_left_if_js = Inches(0.8)
from_top_if_js = Inches(3.3)
from_left_if_py = Inches(5.5)
from_top_if_py = Inches(3.3)
add_if_js = slide7.shapes.add_picture(img_if_js, from_left_if_js, from_top_if_js)
add_if_py = slide7.shapes.add_picture(img_if_py, from_left_if_py, from_top_if_py)


###############################################################################

#Assigning layout to slide 8 - Title and Empty space
slide8_layout = prs.slide_layouts[5]

#Adding slide 8 to presentation
slide8 = prs.slides.add_slide(slide8_layout)

background = slide8.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(253, 239, 180)

#Slide 8 title placeholder
title8 = slide8.shapes.title

#Adding text to slide 8 title
title8.text = "Are you ready to RUUUMBLE?! - Python vs JavaScript"

#Adding text box to appear in the middle of the slide, below title
left8 = Inches(3.8)
top8 = Inches(1.4)
width8 = height8 = Inches(1)
txBox8 = slide8.shapes.add_textbox(left8, top8, width8, height8)
t_in8 = txBox8.text_frame

p8 = t_in8.add_paragraph()
p8.text = "Functions"
p8.font.bold = True
p8.font.size = Pt(30)

#Adding images
img_func_js = "img/functions_js.PNG"
img_func_py = "img/functions_py.PNG"

from_left_func_js = Inches(1.7)
from_top_func_js = Inches(2.3)
from_left_func_py = Inches(2.1)
from_top_func_py = Inches(5)
add_func_js = slide8.shapes.add_picture(img_func_js, from_left_func_js, from_top_func_js)
add_func_py = slide8.shapes.add_picture(img_func_py, from_left_func_py, from_top_func_py)


###############################################################################


#Assign layout type to final slide - Title and blank 
final_slide_layout = prs.slide_layouts[5]

#Add slide 1 to presentation
final_slide = prs.slides.add_slide(final_slide_layout)

background = final_slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(180, 206, 228)

#Main title placeholder
final_title = final_slide.shapes.title

#Adding text to main title placeholder
final_title.text = "And now, for something completely different..."

#Adding images
img_meme1 = "img/python_meme1.jpg"
img_meme2 = "img/python_meme2.jpg"
img_meme3 = "img/python_meme3.jpg"
img_meme4 = "img/python_meme4.png"

from_left_meme1 = Inches(0.5)
from_top_meme1 = Inches(1.8)

from_left_meme2 = Inches(6.5)
from_top_meme2 = Inches(1.8)

from_left_meme3 = Inches(4.3)
from_top_meme3 = Inches(3)

from_left_meme4 = Inches(6.5)
from_top_meme4 = Inches(4.5)

add_meme1 = final_slide.shapes.add_picture(img_meme1, from_left_meme1, from_top_meme1)
add_meme2 = final_slide.shapes.add_picture(img_meme2, from_left_meme2, from_top_meme2)
add_meme3 = final_slide.shapes.add_picture(img_meme3, from_left_meme3, from_top_meme3)
add_meme4 = final_slide.shapes.add_picture(img_meme4, from_left_meme4, from_top_meme4)


###############################################################################


#Saves above logic to .pptx file called CN_PPT_Sofia
prs.save('CN_PPT_Sofia.pptx')